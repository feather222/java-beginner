from docx import Document
from docx.shared import Cm, Pt
from docx.document import Document as Doc
#创建代表Word的doc对象
document = Document() #type:Doc
#添加大标题
document.add_heading('快快乐乐学python',0)
#添加段落
p = document.add_paragraph('python是一门非常流行的语言,它')
#形成新段落内容
run = p.add_run('简单')
#加粗
run.bold = True
run.font.size = Pt(18)
p.add_run('而且')
run = p.add_run('优雅')
run.font.size = Pt(18)
run.underline = True
p.add_run('。')
#添加一级标题
document.add_heading('Heading,level 1', level=1)
#添加带样式的段落
document.add_paragraph('Intense quote',style='Intense Quote')
#添加无序列表(添加列表项目符号)
document.add_paragraph(
    'first item in unordered list', style='List Bullet'
)
document.add_paragraph(
    'second item in unordered list', style='List Bullet'
)
#添加有序列表(添加列表编号)
document.add_paragraph(
    'first item in ordered list', style='List Number'
)
document.add_paragraph(
    'second item in ordered list', style='List Number'
)
#添加照片
document.add_picture(r"D:\图片\屏幕截图\屏幕截图 2026-07-04 184305.png", width = Cm(5.5))
#添加分节符
document.add_section()
records = (
    ('骆昊','男','1995-5-5'),
    ('one','male','2000-1-1')
)
#添加表格
table = document.add_table(rows = 1, cols = 3)
#表格样式
table.style = 'Dark List'
#获取第一行所有单元格
hdr_cells = table.rows[0].cells
#填入内容
hdr_cells[0].text = 'name'
hdr_cells[1].text = 'sex'
hdr_cells[2].text = 'birthday'
#为表格添加行
#元组解包
for name,sex,birthday in records:
    row_cells = table.add_row().cells
    row_cells[0].text = name
    row_cells[1].text = sex
    row_cells[2].text = birthday
#添加分页符
document.add_page_break()
#保存文档
document.save('demo.docx')

from docx import Document
from docx.document import Document as Doc
doc = Document('D:\PyCharm\PythonProjects\demo.docx')
#通过遍历enumerate输出内容,no在这里相当于段落序号
for no, p in enumerate(doc.paragraphs):
    print(no, p.text)

from docx import Document
from docx.document import Document as Doc
#sdate入职时间，edate离职时间
employees = [
    {
        'name':'骆昊',
        'id':'1234',
        'sdate':'2012年2月29日',
        'edate':'2014年3月24日',
        'department':'产品研发',
        'position':'架构师',
        'company':'成都华为技术有限公司'
    },
    {
        'name':'王大锤',
        'id':'1234',
        'sdate':'2019年4月4日',
        'edate':'2021年4月30日',
        'department':'产品研发',
        'position':'Python开发工程师',
        'company':'成都谷道科技有限公司',
    },
    {
        'name':'李元芳',
        'id':'1234',
        'sdate':'2014年5月10日',
        'edate':'2019年3月5日',
        'department':'产品研发',
        'position':'Java开发工程师',
        'company':'同城企业管理集团有限公司'
    },
]
#对列表进行循环遍历，批量生产Word文档
for emp_dict in employees:
#读取已经制作好的辞职信模板
    doc = Document('D:\文档\辞职信.docx')    #Doc
#循环遍历所有段落，找出占位符
    for p in doc.paragraphs:
        if '{' not in p.text:
            continue
        for run in p.runs:
            if '{'not in run.text:
                continue
            #将占位符{换成实际内容
            start, end = run.text.find('{'), run.text.find('}')
            key, place_holder = run.text[start + 1:end], run.text[start:end+1]
            run.text = run.text.replace(place_holder, emp_dict[key])
#保存文档
    doc.save(f'{emp_dict["name"]}离职证明.docx')

#python对powerpoint的操作
from pptx import Presentation
#创建幻灯片对象
pres = Presentation()
#选择母版添加一页
#取出index为0的版式（标题幻灯片）
title_slide_layout = pres.slide_layouts[0]
#添加0版式的幻灯片
slide = pres.slides.add_slide(title_slide_layout)
#获取标题栏（标题栏包括文本框、图片、图表等）
title = slide.shapes.title
#获取副标题
subtitle = slide.placeholders[1]
#编辑标题和副标题
title.text = 'Welcome to python'
subtitle.text = 'Life is short, I use Python'
#选择母版添加一页，添加1版式的幻灯片
bullet_slide_layout = pres.slide_layouts[1]
slide = pres.slides.add_slide(bullet_slide_layout)
#获取页面上所有文本框
shapes = slide.shapes
#获取标题文本框和正文文本框
title_shape = shapes.title
body_shape = shapes.placeholders[1]
#编辑内容
title_shape.text = 'Introduction'
#编辑正文文本框内容
tf = body_shape.text_frame
tf.text = 'History of Python'
#添加一个一级标题
p = tf.add_paragraph()
#添加内容
p.text = "X'max 1989"
p.level = 1
#添加一个二级标题
p = tf.add_paragraph()
p.text = 'Guido began to write interpreter for Python'
p.level = 2
#保存幻灯片
pres.save('test.pptx')
